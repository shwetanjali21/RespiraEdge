import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize 16:9 Widescreen Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- EY CORPORATE THEME PALETTE ---
EY_YELLOW = RGBColor(255, 230, 0)     # Primary Accent (#FFE600)
EY_CHARCOAL = RGBColor(46, 46, 56)    # Dark Slide Background (#2E2E38)
TEXT_WHITE = RGBColor(245, 245, 247)   # High-Contrast Body Text
TEXT_MUTED = RGBColor(180, 180, 190)   # Muted Subtitles & Descriptions

# Path to your circuit image in the workspace folder
IMAGE_PATH = "Screenshot 2026-08-29 144630.png"

def set_slide_background(slide):
    """Fills slide background with EY Dark Charcoal."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = EY_CHARCOAL

def add_header(slide, title_text):
    """Creates a standardized slide title block."""
    set_slide_background(slide)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = EY_YELLOW

blank_layout = prs.slide_layouts[6]

# =============================================================================
# SLIDE 1: Title Slide
# =============================================================================
slide1 = prs.slides.add_slide(blank_layout)
set_slide_background(slide1)

box1 = slide1.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(3.8))
tf1 = box1.text_frame
tf1.word_wrap = True

p1 = tf1.paragraphs[0]
p1.text = "RespiraEdge"
p1.font.size = Pt(48)
p1.font.bold = True
p1.font.color.rgb = EY_YELLOW

p2 = tf1.add_paragraph()
p2.text = "On-Device Edge AI for Autonomous Respiratory Distress & Apnea Detection"
p2.font.size = Pt(22)
p2.font.color.rgb = TEXT_WHITE
p2.space_before = Pt(12)

p3 = tf1.add_paragraph()
p3.text = "Presenter: Shwetanjali Gautam  |  Domain: Embedded Systems, Edge AI & Medical IoT"
p3.font.size = Pt(14)
p3.font.color.rgb = TEXT_MUTED
p3.space_before = Pt(36)

# =============================================================================
# SLIDE 2: Executive Summary (Cloud vs. Edge)
# =============================================================================
slide2 = prs.slides.add_slide(blank_layout)
add_header(slide2, "Executive Summary: Cloud vs. Edge AI")

box2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
tf2 = box2.text_frame
tf2.word_wrap = True

p = tf2.paragraphs[0]
p.text = "The Problem with Cloud-Based Health IoT:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = EY_YELLOW

bullets_prob = [
    "High Latency: Round-trip delays prevent instant response during acute respiratory events.",
    "Privacy Risks: Transmitting sensitive patient bio-signals over public networks.",
    "Connectivity Dependency: System failure during network outages or poor reception."
]
for b in bullets_prob:
    p = tf2.add_paragraph()
    p.text = "• " + b
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

p = tf2.add_paragraph()
p.text = "\nThe RespiraEdge Solution:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = EY_YELLOW
p.space_before = Pt(16)

bullets_sol = [
    "100% On-Device Neural Network running directly on an ESP32 microcontroller.",
    "Sub-millisecond latency (200 to 400 microseconds execution time).",
    "Autonomous local closed-loop actuation with zero cloud dependency."
]
for b in bullets_sol:
    p = tf2.add_paragraph()
    p.text = "• " + b
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

# =============================================================================
# SLIDE 3: Hardware Pin Mapping & Circuit Layout (With Image)
# =============================================================================
slide3 = prs.slides.add_slide(blank_layout)
add_header(slide3, "Hardware Circuit Layout & Pin Connections")

# Table on left half
rows, cols = 6, 3
left, top, width, height = Inches(0.8), Inches(1.6), Inches(6.8), Inches(5.2)
table_shape = slide3.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

headers = ["Component", "Pin", "Technical Function"]
data = [
    ["DHT22 Sensor", "GPIO 4", "Monitors ambient humidity (%) & temp (°C)"],
    ["Potentiometer", "GPIO 36", "Simulates continuous audio RMS energy stream"],
    ["Relay Module", "GPIO 25", "Actuates emergency humidifier / valve"],
    ["Piezo Buzzer", "GPIO 26", "Audible alarm indicator for distress states"],
    ["Red LED", "GPIO 33", "Visual alert indicator"]
]

for col_idx, text in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = EY_YELLOW
    for p in cell.text_frame.paragraphs:
        p.font.color.rgb = EY_CHARCOAL
        p.font.bold = True
        p.font.size = Pt(12)

for row_idx, row_data in enumerate(data):
    for col_idx, text in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(60, 60, 72)
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = TEXT_WHITE
            p.font.size = Pt(11)

# Embed Image on right half if available
if os.path.exists(IMAGE_PATH):
    slide3.shapes.add_picture(IMAGE_PATH, Inches(7.9), Inches(1.6), width=Inches(4.6))
else:
    # Placeholder box if image path is incorrect
    img_box = slide3.shapes.add_textbox(Inches(7.9), Inches(1.6), Inches(4.6), Inches(5.2))
    img_tf = img_box.text_frame
    img_tf.word_wrap = True
    p = img_tf.paragraphs[0]
    p.text = f"[Circuit Screenshot Placeholder]\n\nAdd '{IMAGE_PATH}' to your project folder to display it here."
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED

# =============================================================================
# SLIDE 4: Machine Learning Engine & C Code Architecture
# =============================================================================
slide4 = prs.slides.add_slide(blank_layout)
add_header(slide4, "On-Device Neural Network Engine")

box4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
tf4 = box4.text_frame
tf4.word_wrap = True

bullets_ml = [
    "Offline Model Training (RespiraEdge.ipynb): Multi-layer feedforward ANN trained in Python.",
    "C Weight Export (model_weights.h): Weights & biases converted to static C floating-point matrices.",
    "Embedded Inference (sketch.ino): Execution via forward propagation: y = ReLU(W · x + b).",
    "Network Topology: 4 Feature Inputs → 16 Neurons (ReLU) → 8 Neurons (ReLU) → 3 Classes.",
    "Class Outputs: [0: Normal State]  |  [1: Wheezing Distress]  |  [2: Apnea Event]"
]
for b in bullets_ml:
    p = tf4.add_paragraph()
    p.text = "• " + b
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(14)

# =============================================================================
# SLIDE 5: Performance Metrics & Future Scope
# =============================================================================
slide5 = prs.slides.add_slide(blank_layout)
add_header(slide5, "Performance Metrics & Future Roadmap")

box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
tf5 = box5.text_frame
tf5.word_wrap = True

p = tf5.paragraphs[0]
p.text = "Key System Metrics:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = EY_YELLOW

achievements = [
    "Inference Speed: Sub-millisecond execution time (200-400 µs per prediction).",
    "Data Privacy: Bio-signals processed 100% locally with zero transmission.",
    "Hardware Overhead: Cost-effective design operating on standard ESP32 MCU."
]
for a in achievements:
    p = tf5.add_paragraph()
    p.text = "• " + a
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(5)

p = tf5.add_paragraph()
p.text = "\nFuture Scope:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = EY_YELLOW
p.space_before = Pt(16)

future = [
    "Digital Microphone Integration: Replace potentiometer with INMP441 MEMS mic via I2S.",
    "INT8 Quantization: Convert float weights to fixed-point integers for faster execution.",
    "BLE Alerts: Broadcast Bluetooth Low Energy signals to notify local health responders."
]
for f in future:
    p = tf5.add_paragraph()
    p.text = "• " + f
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(5)

# Save output
prs.save("RespiraEdge_Presentation.pptx")
print("Successfully generated EY-themed RespiraEdge_Presentation.pptx with image support!")